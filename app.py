from flask import Flask, render_template, request
import pandas as pd
from pptx import Presentation
from googletrans import Translator
import os

app = Flask(__name__)

@app.route('/', methods=['GET', 'POST'])
def index():
    if request.method == 'POST':
        target_language = request.form['target_language']
        uploaded_file = request.files['file']

        if uploaded_file.filename.endswith(('.xls', '.xlsx')):
            # Handle Excel file
            df = pd.read_excel(uploaded_file)
            translated_df = translate_dataframe(df, target_language)
            translated_file_path = save_translated_excel(translated_df)

        elif uploaded_file.filename.endswith('.pptx'):
            # Handle PowerPoint file
            prs = Presentation(uploaded_file)
            translated_prs = translate_pptx(prs, target_language)
            translated_file_path = save_translated_pptx(translated_prs)

        else:
            return render_template('index.html', error='Unsupported file format')

        return render_template('index.html', translated_file=translated_file_path)

    return render_template('index.html')

def translate_dataframe(df, target_language):
    # Create a Translator object
    translator = Translator()

    # Translate the header (row 1)
    translated_header = df.columns.map(lambda col: translator.translate(str(col), dest=target_language).text)

    # Function to translate a single cell
    def translate_cell(cell):
        return translator.translate(str(cell), dest=target_language).text

    # Apply translation to each cell in the DataFrame
    translated_df = df.applymap(translate_cell)

    # Replace the header with the translated header
    translated_df.columns = translated_header

    return translated_df

def save_translated_excel(translated_df):
    # Output file path for the translated Excel file
    translated_file_path = 'static/translated_output.xlsx'

    # Save the translated DataFrame to an Excel file
    translated_df.to_excel(translated_file_path, index=False)

    return translated_file_path

def translate_pptx(prs, target_language):
    # Create a Translator object
    translator = Translator()

    # Function to translate text within a shape
    def translate_shape_text(shape):
        if hasattr(shape, 'text'):
            shape.text = translator.translate(shape.text, dest=target_language).text

    # Apply translation to text in each shape of each slide
    for slide in prs.slides:
        for shape in slide.shapes:
            translate_shape_text(shape)

    return prs

def save_translated_pptx(translated_prs):
    # Output file path for the translated PowerPoint file
    translated_file_path = 'static/translated_output.pptx'

    # Save the translated PowerPoint presentation to a file
    translated_prs.save(translated_file_path)

    return translated_file_path

if __name__ == '__main__':
    app.run(debug=True)
